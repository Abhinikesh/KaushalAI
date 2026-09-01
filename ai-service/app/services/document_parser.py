from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx"}


def parse_pdf(file_bytes: bytes) -> str:
    import fitz  # PyMuPDF

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n".join(pages)
    except Exception as exc:
        raise ValueError(f"Failed to parse PDF: {exc}") from exc


def parse_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except Exception as exc:
        raise ValueError(f"Failed to parse PPTX: {exc}") from exc


def parse_docx(file_bytes: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as exc:
        raise ValueError(f"Failed to parse DOCX: {exc}") from exc


def parse_document(filename: str, file_bytes: bytes) -> str:
    """Dispatch to the correct parser based on file extension."""
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    parsers = {".pdf": parse_pdf, ".pptx": parse_pptx, ".docx": parse_docx}
    text = parsers[ext](file_bytes)

    if not text.strip():
        raise ValueError(f"Document '{filename}' produced no extractable text. It may be image-only or corrupt.")

    logger.info("Parsed '%s' (%s): %d chars extracted", filename, ext, len(text))
    return text
